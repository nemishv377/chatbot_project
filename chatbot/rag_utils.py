import hashlib
import io
import mimetypes
import os
import uuid
from typing import List
from typing import Tuple

import chromadb

# File parsers
import fitz  # PyMuPDF (PDF)
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation as PptxPresentation
from sentence_transformers import SentenceTransformer

# OCR
try:
    import pytesseract

    OCR_AVAILABLE = True
except Exception:
    OCR_AVAILABLE = False

# -------------------------
# Persistent Chroma setup
# -------------------------
CHROMA_PATH = "./chroma_db"
chroma_client = chromadb.PersistentClient(path=CHROMA_PATH)
collection = chroma_client.get_or_create_collection(name="my_docs")

# Embedding model
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")


# -------------------------
# Helpers
# -------------------------
def sha1(text: str) -> str:
    return hashlib.sha1(text.encode("utf-8", errors="ignore")).hexdigest()


# def split_text(text: str, chunk_size: int = 300, overlap: int = 50) -> List[str]:
def split_text(text: str, chunk_size: int = 800, overlap: int = 120) -> List[str]:
    """Simple overlap chunking that prefers splitting on paragraph boundaries."""
    text = text.replace("\r", "")
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    chunks, buf = [], ""
    for p in paragraphs:
        if len(buf) + len(p) + 1 <= chunk_size:
            buf = f"{buf}\n{p}" if buf else p
        else:
            if buf:
                chunks.append(buf)
            buf = p
    if buf:
        chunks.append(buf)

    # If paragraphs were huge, fall back to raw slicing with overlap
    final_chunks = []
    for c in chunks:
        if len(c) <= chunk_size:
            final_chunks.append(c)
        else:
            start = 0
            while start < len(c):
                end = start + chunk_size
                final_chunks.append(c[start:end])
                start += chunk_size - overlap
    return final_chunks


# -------------------------
# Extraction per type
# -------------------------
def extract_text_from_pdf(path: str) -> str:
    doc = fitz.open(path)
    pages = []
    for i in range(len(doc)):
        pages.append(doc[i].get_text("text"))
    return "\n".join(pages)


def extract_text_from_docx(path: str) -> str:
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_pptx(path: str) -> str:
    prs = PptxPresentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)


def extract_text_from_excel(path: str) -> str:
    # Handles xlsx/xls
    dfs = pd.read_excel(path, sheet_name=None)
    out = []
    for sheet_name, df in dfs.items():
        out.append(f"# Sheet: {sheet_name}\n{df.to_string(index=False)}")
    return "\n\n".join(out)


def extract_text_from_csv(path: str) -> str:
    df = pd.read_csv(path)
    return df.to_string(index=False)


def extract_text_from_html(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "lxml")
    return soup.get_text(separator="\n", strip=True)


def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text_from_image(path: str) -> str:
    if not OCR_AVAILABLE:
        return ""
    img = Image.open(path)
    # Optional: convert to grayscale to boost OCR a bit
    try:
        img = img.convert("L")
    except Exception:
        pass
    return pytesseract.image_to_string(img)


# -------------------------
# Router
# -------------------------
EXT_MAP = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".pptx": extract_text_from_pptx,
    ".xlsx": extract_text_from_excel,
    ".xls": extract_text_from_excel,
    ".csv": extract_text_from_csv,
    ".html": extract_text_from_html,
    ".htm": extract_text_from_html,
    ".txt": extract_text_from_txt,
    ".md": extract_text_from_txt,
    ".jpg": extract_text_from_image,
    ".jpeg": extract_text_from_image,
    ".png": extract_text_from_image,
    ".webp": extract_text_from_image,
}


def extract_text_generic(path: str) -> Tuple[str, str]:
    """
    Return (text, used_extractor).
    """
    ext = os.path.splitext(path)[1].lower()
    func = EXT_MAP.get(ext)
    if func is None:
        # Fallback by MIME/extension
        if ext in [".gif", ".tif", ".tiff"]:
            func = extract_text_from_image
        else:
            func = extract_text_from_txt  # attempt as text
    text = ""
    try:
        text = func(path)
    except Exception:
        text = ""
    return text, func.__name__


# -------------------------
# Indexing / Retrieval
# -------------------------
def add_document_text(full_text: str, source_name: str) -> int:
    """
    Split, embed, and add text to ChromaDB.
    Returns number of chunks indexed.
    """
    chunks = [c for c in split_text(full_text) if len(c.strip()) > 0]
    if not chunks:
        return 0

    embeddings = embedding_model.encode(chunks, convert_to_numpy=True).tolist()
    print("================")
    ids = [
        f"{sha1(source_name)}_{i}_{uuid.uuid4().hex[:8]}" for i in range(len(chunks))
    ]

    metadatas = [{"source": source_name, "chunk": i} for i in range(len(chunks))]

    # Note: Chroma will error on duplicate ids; for repeated ingests, you may append random suffix.
    collection.add(
        documents=chunks, embeddings=embeddings, ids=ids, metadatas=metadatas
    )
    return len(chunks)


def ingest_file_to_chroma(path: str, source_name: str) -> Tuple[int, str]:
    """
    Extract text from the file and index it.
    Returns (num_chunks, extractor_name)
    """
    text, extractor = extract_text_generic(path)
    if not text or len(text.strip()) == 0:
        return 0, extractor
    num = add_document_text(text, source_name)
    return num, extractor


def get_relevant_chunks(query: str, top_k: int = 5) -> str:
    """Retrieve most relevant doc chunks for a query."""
    q_emb = embedding_model.encode(query).tolist()
    results = collection.query(query_embeddings=[q_emb], n_results=top_k)
    return (
        "\n".join(results["documents"][0])
        if results and results.get("documents")
        else ""
    )
